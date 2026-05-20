# тут у нас конвертер файлов, я не уверен что это всё работает но попробуем
# поддерживает: документы, картинки, аудио/видео, всякое говно
# ffmpeg клади в ./ffmpeg/ffmpeg иначе облом

import csv  # тут импорты, я даун
import json
import os
import re
import shutil
import subprocess
import tempfile
import xml.etree.ElementTree as ET  # для xml, хз зачем так называется
from pathlib import Path
from typing import Optional, Union

from config import get_ffmpeg, get_libreoffice  # это из моего файла конфига

# ── какие форматы мы умеем ── (я сам не знаю все но звучит круто)
IMAGE_FMTS = {"png", "jpg", "jpeg", "webp", "bmp", "gif", "tiff", "tif", "ico"}
VIDEO_FMTS = {"mp4", "avi", "mkv", "mov", "webm", "flv", "m4v", "wmv", "3gp", "ts", "vob", "mpeg", "mpg"}
AUDIO_FMTS = {"mp3", "wav", "flac", "aac", "ogg", "m4a", "wma", "opus", "ac3", "mp2"}
MEDIA_FMTS = VIDEO_FMTS | AUDIO_FMTS  # объединяю два сета, умный же

# всё что можем взять и всё что можем выдать
SUPPORTED_IN  = {"xml", "json", "csv", "yaml", "yml", "html", "docx", "doc",
                 "pptx", "ppt", "xlsx", "xls", "pdf", "txt", "md", "markdown"} | IMAGE_FMTS | MEDIA_FMTS
SUPPORTED_OUT = {"xml", "json", "csv", "yaml", "yml", "html", "docx", "pdf",
                 "txt", "md"} | IMAGE_FMTS | MEDIA_FMTS


# вот это главная функция, она всё решает
def convert(src: Union[str, Path], dst_fmt: str,
            dst: Optional[Union[str, Path]] = None) -> Path:
    src     = Path(src)  # делаю нормальным путём
    src_fmt = src.suffix.lstrip(".").lower()  # берём расширение
    dst_fmt = dst_fmt.lstrip(".").lower()

    if dst is None:
        dst = src.with_suffix(f".{dst_fmt}")  # если не сказали куда — рядом кладём
    dst = Path(dst)
    dst.parent.mkdir(parents=True, exist_ok=True)  # создаём папку если нет

    # если медиа — сразу в ffmpeg, он умнее меня
    if src_fmt in MEDIA_FMTS or dst_fmt in MEDIA_FMTS:
        return _ffmpeg(src, dst)

    # картинки через пилоу (это библиотека, не птица)
    if src_fmt in IMAGE_FMTS or dst_fmt in IMAGE_FMTS:
        return _image(src, dst, dst_fmt)

    # ищем нужную функцию в таблице снизу
    key     = (src_fmt, dst_fmt)
    handler = _DISPATCH.get(key)
    if handler is None:
        # может либреоффис вывезет?
        office_in  = {"doc", "docx", "ppt", "pptx", "xls", "xlsx", "odt", "ods", "odp"}
        office_out = {"pdf", "html", "txt", "png", "docx"}
        if src_fmt in office_in and dst_fmt in office_out:
            return _libreoffice(src, dst, dst_fmt)
        raise ValueError(f"Не умею делать {src_fmt} → {dst_fmt}, сорри")

    handler(src, dst)  # вызываем функцию конвертации
    return dst


# ── ffmpeg часть ── (это не я написал, это он сам умный)

def _ffmpeg(src: Path, dst: Path) -> Path:
    ff = get_ffmpeg()
    if not ff:
        raise RuntimeError(
            "FFmpeg не найден! Скачай и положи в папку ./ffmpeg/ffmpeg\n"
            "Качать тут: https://ffmpeg.org/download.html"
        )
    # запускаем ffmpeg и молимся
    result = subprocess.run(
        [ff, "-y", "-i", str(src), str(dst)],
        capture_output=True, text=True
    )
    if result.returncode != 0:
        raise RuntimeError(f"FFmpeg упал с ошибкой:\n{result.stderr[-800:]}")
    return dst


# ── либреоффис для конвертации офисных файлов ── (тяжёлая артиллерия)

def _libreoffice(src: Path, dst: Path, dst_fmt: str) -> Path:
    lo = get_libreoffice()
    if not lo:
        raise RuntimeError("LibreOffice не стоит. Ставь: apt install libreoffice")
    with tempfile.TemporaryDirectory() as tmp:
        # он конвертит в папку, потом ищем результат
        subprocess.run(
            [lo, "--headless", "--convert-to", dst_fmt, "--outdir", tmp, str(src)],
            capture_output=True, check=True
        )
        converted = next(Path(tmp).glob(f"*.{dst_fmt}"), None)
        if not converted:
            raise RuntimeError("LibreOffice ничего не создал, странно")
        shutil.move(str(converted), str(dst))
    return dst


# ── картинки через Pillow ──

def _image(src: Path, dst: Path, fmt: str) -> Path:
    try:
        from PIL import Image  # если нет — говорим что ставить
    except ImportError:
        raise ImportError("pip install Pillow")
    # некоторые форматы требуют RGB, без альфа канала
    fmt_map = {"jpg": "JPEG", "jpeg": "JPEG", "tif": "TIFF"}
    pil_fmt = fmt_map.get(fmt, fmt.upper())
    with Image.open(src) as img:
        if pil_fmt == "JPEG" and img.mode in ("RGBA", "P"):
            img = img.convert("RGB")  # убираем прозрачность иначе JPEG плачет
        img.save(dst, pil_fmt)
    return dst


# ── xml → всякое ── (xml это боль но справимся)

def _elem_to_dict(elem: ET.Element):
    # рекурсивно превращаем xml в питоновский словарь
    d = {}
    for child in elem:
        val = _elem_to_dict(child)
        if child.tag in d:
            if not isinstance(d[child.tag], list):
                d[child.tag] = [d[child.tag]]
            d[child.tag].append(val)
        else:
            d[child.tag] = val
    if elem.attrib:
        d["@attr"] = elem.attrib
    if elem.text and elem.text.strip():
        return elem.text.strip() if not d else {**d, "#text": elem.text.strip()}
    return d or (elem.text or "")


def _dict_to_elem(tag: str, data) -> ET.Element:
    # обратно, словарь → xml элемент
    elem = ET.Element(tag)
    if isinstance(data, dict):
        for k, v in data.items():
            if k == "@attr":
                elem.attrib.update(v)
            elif k == "#text":
                elem.text = str(v)
            elif isinstance(v, list):
                for item in v:
                    elem.append(_dict_to_elem(k, item))
            else:
                elem.append(_dict_to_elem(k, v))
    elif isinstance(data, list):
        for item in data:
            elem.append(_dict_to_elem("item", item))
    else:
        elem.text = str(data)
    return elem


def _xml_to_json(src: Path, dst: Path):
    tree = ET.parse(src)
    root = tree.getroot()
    data = {root.tag: _elem_to_dict(root)}
    dst.write_text(json.dumps(data, ensure_ascii=False, indent=2))


def _xml_to_yaml(src: Path, dst: Path):
    try:
        import yaml
    except ImportError:
        raise ImportError("pip install PyYAML")
    tree  = ET.parse(src)
    root  = tree.getroot()
    data  = {root.tag: _elem_to_dict(root)}
    dst.write_text(yaml.dump(data, allow_unicode=True, default_flow_style=False))


def _xml_to_csv(src: Path, dst: Path):
    # берём всех детей корня и делаем из них строки
    tree    = ET.parse(src)
    root    = tree.getroot()
    rows    = []
    headers: list[str] = []
    for child in root:
        row = {}
        for sub in child:
            if sub.tag not in headers:
                headers.append(sub.tag)
            row[sub.tag] = sub.text or ""
        rows.append(row)
    with open(dst, "w", newline="", encoding="utf-8") as f:
        w = csv.DictWriter(f, fieldnames=headers)
        w.writeheader()
        w.writerows(rows)


def _xml_to_html(src: Path, dst: Path):
    tree = ET.parse(src)
    root = tree.getroot()
    data = {root.tag: _elem_to_dict(root)}
    html = _dict_to_html_table(data)
    dst.write_text(
        f"<!DOCTYPE html><html><head><meta charset='utf-8'><title>{root.tag}</title></head>"
        f"<body style='font-family:sans-serif'>{html}</body></html>"
    )


def _dict_to_html_table(data, depth=0) -> str:
    # рекурсивно делаем таблицу из словаря
    if isinstance(data, dict):
        rows = "".join(
            f"<tr><td><b>{k}</b></td><td>{_dict_to_html_table(v, depth+1)}</td></tr>"
            for k, v in data.items()
        )
        return f"<table border='1' cellpadding='4' style='border-collapse:collapse'>{rows}</table>"
    elif isinstance(data, list):
        items = "".join(f"<li>{_dict_to_html_table(i, depth+1)}</li>" for i in data)
        return f"<ul>{items}</ul>"
    return str(data)


# ── json → всякое ──

def _json_to_xml(src: Path, dst: Path):
    data = json.loads(src.read_text(encoding="utf-8"))
    if isinstance(data, dict) and len(data) == 1:
        tag, val = next(iter(data.items()))
    else:
        tag, val = "root", data
    root = _dict_to_elem(tag, val)
    ET.indent(root, space="  ")
    ET.ElementTree(root).write(dst, encoding="unicode", xml_declaration=True)


def _json_to_csv(src: Path, dst: Path):
    data = json.loads(src.read_text(encoding="utf-8"))
    if isinstance(data, dict):
        data = [data]  # один объект тоже обрабатываем
    if not data or not isinstance(data[0], dict):
        raise ValueError("JSON должен быть массивом объектов или одним объектом")
    headers = list(data[0].keys())
    with open(dst, "w", newline="", encoding="utf-8") as f:
        w = csv.DictWriter(f, fieldnames=headers, extrasaction="ignore")
        w.writeheader()
        w.writerows(data)


def _json_to_yaml(src: Path, dst: Path):
    try:
        import yaml
    except ImportError:
        raise ImportError("pip install PyYAML")
    data = json.loads(src.read_text(encoding="utf-8"))
    dst.write_text(yaml.dump(data, allow_unicode=True, default_flow_style=False))


# ── csv → всякое ──

def _csv_to_json(src: Path, dst: Path):
    with open(src, encoding="utf-8") as f:
        rows = list(csv.DictReader(f))
    dst.write_text(json.dumps(rows, ensure_ascii=False, indent=2))


def _csv_to_xml(src: Path, dst: Path):
    with open(src, encoding="utf-8") as f:
        rows = list(csv.DictReader(f))
    root = ET.Element("data")
    for row in rows:
        item = ET.SubElement(root, "item")
        for k, v in row.items():
            sub      = ET.SubElement(item, k.replace(" ", "_").replace("-", "_"))
            sub.text = v
    ET.indent(root, space="  ")
    ET.ElementTree(root).write(dst, encoding="unicode", xml_declaration=True)


# ── yaml → всякое ──

def _yaml_to_json(src: Path, dst: Path):
    try:
        import yaml
    except ImportError:
        raise ImportError("pip install PyYAML")
    data = yaml.safe_load(src.read_text(encoding="utf-8"))
    dst.write_text(json.dumps(data, ensure_ascii=False, indent=2))


def _yaml_to_xml(src: Path, dst: Path):
    try:
        import yaml
    except ImportError:
        raise ImportError("pip install PyYAML")
    data = yaml.safe_load(src.read_text(encoding="utf-8"))
    if isinstance(data, dict) and len(data) == 1:
        tag, val = next(iter(data.items()))
    else:
        tag, val = "root", data
    root = _dict_to_elem(tag, val)
    ET.indent(root, space="  ")
    ET.ElementTree(root).write(dst, encoding="unicode", xml_declaration=True)


# ── docx ── (word файлы, мой любимый ужас)

def _docx_to_txt(src: Path, dst: Path):
    try:
        from docx import Document
    except ImportError:
        raise ImportError("pip install python-docx")
    doc = Document(src)
    dst.write_text("\n".join(p.text for p in doc.paragraphs), encoding="utf-8")


def _docx_to_html(src: Path, dst: Path):
    try:
        from docx import Document
    except ImportError:
        raise ImportError("pip install python-docx")
    doc   = Document(src)
    parts = ["<!DOCTYPE html><html><head><meta charset='utf-8'>"
             "<style>body{font-family:sans-serif;max-width:800px;margin:40px auto;line-height:1.6}</style>"
             "</head><body>"]
    for p in doc.paragraphs:
        style = p.style.name.lower()
        if style.startswith("heading"):
            lvl = style.split()[-1] if style.split()[-1].isdigit() else "1"
            parts.append(f"<h{lvl}>{p.text}</h{lvl}>")
        elif p.text.strip():
            parts.append(f"<p>{p.text}</p>")
    parts.append("</body></html>")
    dst.write_text("\n".join(parts), encoding="utf-8")


def _docx_to_md(src: Path, dst: Path):
    try:
        from docx import Document
    except ImportError:
        raise ImportError("pip install python-docx")
    doc   = Document(src)
    lines = []
    for p in doc.paragraphs:
        style = p.style.name.lower()
        if style.startswith("heading"):
            lvl = int(style.split()[-1]) if style.split()[-1].isdigit() else 1
            lines.append(f"{'#' * lvl} {p.text}")
        else:
            lines.append(p.text)
    dst.write_text("\n\n".join(lines), encoding="utf-8")


def _docx_to_pdf(src: Path, dst: Path):
    # через либреоффис, другого нормального варианта на линуксе нет
    return _libreoffice(src, dst, "pdf")


# ── pptx ── (слайды, тоже боль)

def _pptx_to_txt(src: Path, dst: Path):
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("pip install python-pptx")
    prs   = Presentation(src)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"=== Слайд {i} ===")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text)
    dst.write_text("\n".join(lines), encoding="utf-8")


def _pptx_to_html(src: Path, dst: Path):
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("pip install python-pptx")
    prs   = Presentation(src)
    parts = ["<!DOCTYPE html><html><head><meta charset='utf-8'>"
             "<style>section{border:1px solid #ccc;margin:20px;padding:20px;border-radius:8px}"
             "body{font-family:sans-serif;max-width:900px;margin:40px auto}</style></head><body>"]
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"<section><h2>Слайд {i}</h2>")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(f"<p>{shape.text}</p>")
        parts.append("</section>")
    parts.append("</body></html>")
    dst.write_text("\n".join(parts), encoding="utf-8")


def _pptx_to_pdf(src: Path, dst: Path):
    return _libreoffice(src, dst, "pdf")


# ── xlsx ── (экселька, боже мой)

def _xlsx_to_csv(src: Path, dst: Path):
    try:
        import openpyxl
    except ImportError:
        raise ImportError("pip install openpyxl")
    wb = openpyxl.load_workbook(src, read_only=True, data_only=True)
    ws = wb.active  # берём первый лист
    with open(dst, "w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        for row in ws.iter_rows(values_only=True):
            w.writerow([str(c) if c is not None else "" for c in row])


def _xlsx_to_json(src: Path, dst: Path):
    try:
        import openpyxl
    except ImportError:
        raise ImportError("pip install openpyxl")
    wb   = openpyxl.load_workbook(src, read_only=True, data_only=True)
    ws   = wb.active
    rows = list(ws.iter_rows(values_only=True))
    if not rows:
        dst.write_text("[]")
        return
    headers = [str(h) if h is not None else f"col{i}" for i, h in enumerate(rows[0])]
    data    = [dict(zip(headers, [str(c) if c is not None else "" for c in row]))
               for row in rows[1:]]
    dst.write_text(json.dumps(data, ensure_ascii=False, indent=2))


def _xlsx_to_html(src: Path, dst: Path):
    try:
        import openpyxl
    except ImportError:
        raise ImportError("pip install openpyxl")
    wb   = openpyxl.load_workbook(src, read_only=True, data_only=True)
    ws   = wb.active
    rows = list(ws.iter_rows(values_only=True))
    # красивая таблица с заголовками
    parts = ["<!DOCTYPE html><html><head><meta charset='utf-8'>"
             "<style>table{border-collapse:collapse;width:100%}th,td{border:1px solid #ddd;padding:8px}"
             "th{background:#f2f2f2}body{font-family:sans-serif;padding:20px}</style></head><body>",
             "<table>"]
    for i, row in enumerate(rows):
        tag   = "th" if i == 0 else "td"
        cells = "".join(f"<{tag}>{c if c is not None else ''}</{tag}>" for c in row)
        parts.append(f"<tr>{cells}</tr>")
    parts.append("</table></body></html>")
    dst.write_text("\n".join(parts), encoding="utf-8")


# ── pdf ── (читать умеем, писать через reportlab)

def _pdf_to_txt(src: Path, dst: Path):
    try:
        import pdfplumber  # первый вариант
        with pdfplumber.open(src) as pdf:
            text = "\n\n".join(p.extract_text() or "" for p in pdf.pages)
        dst.write_text(text, encoding="utf-8")
        return
    except ImportError:
        pass
    try:
        import pypdf  # второй вариант
        reader = pypdf.PdfReader(str(src))
        text   = "\n\n".join(p.extract_text() or "" for p in reader.pages)
        dst.write_text(text, encoding="utf-8")
        return
    except ImportError:
        pass
    raise ImportError("нужно одно из: pip install pdfplumber  ИЛИ  pip install pypdf")


def _pdf_to_html(src: Path, dst: Path):
    # сначала в текст, потом в html
    tmp = src.with_suffix(".tmp.txt")
    _pdf_to_txt(src, tmp)
    content = tmp.read_text(encoding="utf-8")
    tmp.unlink(missing_ok=True)
    paragraphs = "\n".join(f"<p>{line}</p>" for line in content.splitlines() if line.strip())
    dst.write_text(
        f"<!DOCTYPE html><html><head><meta charset='utf-8'>"
        f"<style>body{{font-family:sans-serif;max-width:800px;margin:40px auto;line-height:1.6}}</style>"
        f"</head><body>{paragraphs}</body></html>"
    )


# ── txt/html/markdown разные штуки ──

def _txt_to_html(src: Path, dst: Path):
    text       = src.read_text(encoding="utf-8")
    paragraphs = "\n".join(f"<p>{line}</p>" for line in text.splitlines() if line.strip())
    dst.write_text(
        f"<!DOCTYPE html><html><head><meta charset='utf-8'></head><body>{paragraphs}</body></html>"
    )


def _txt_to_pdf(src: Path, dst: Path):
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import SimpleDocTemplate, Paragraph
    except ImportError:
        raise ImportError("pip install reportlab")
    text  = src.read_text(encoding="utf-8")
    doc   = SimpleDocTemplate(str(dst), pagesize=A4)
    styles = getSampleStyleSheet()
    # каждая строка — параграф, пустые строки — отступ
    story = [Paragraph(line or "&nbsp;", styles["Normal"]) for line in text.splitlines()]
    doc.build(story)


def _md_to_html(src: Path, dst: Path):
    try:
        import markdown as mdlib  # библиотека markdown, не путать с форматом
    except ImportError:
        raise ImportError("pip install Markdown")
    text     = src.read_text(encoding="utf-8")
    html_body = mdlib.markdown(text, extensions=["tables", "fenced_code"])
    dst.write_text(
        f"<!DOCTYPE html><html><head><meta charset='utf-8'>"
        f"<style>body{{font-family:sans-serif;max-width:900px;margin:40px auto;line-height:1.6}}"
        f"code{{background:#f4f4f4;padding:2px 6px;border-radius:4px}}"
        f"pre{{background:#f4f4f4;padding:16px;border-radius:8px;overflow:auto}}</style>"
        f"</head><body>{html_body}</body></html>"
    )


def _html_to_txt(src: Path, dst: Path):
    text = src.read_text(encoding="utf-8")
    try:
        from bs4 import BeautifulSoup
        text = BeautifulSoup(text, "html.parser").get_text("\n")
    except ImportError:
        # если нет bs4 — просто убираем теги регуляркой, грубо но работает
        text = re.sub(r"<[^>]+>", "", text)
    dst.write_text(text, encoding="utf-8")


# ── таблица диспетчеризации ── (тут решается кто что конвертит)

_DISPATCH = {
    # xml во всё
    ("xml",  "json"):    _xml_to_json,
    ("xml",  "yaml"):    _xml_to_yaml,
    ("xml",  "yml"):     _xml_to_yaml,
    ("xml",  "csv"):     _xml_to_csv,
    ("xml",  "html"):    _xml_to_html,
    # json во всё
    ("json", "xml"):     _json_to_xml,
    ("json", "csv"):     _json_to_csv,
    ("json", "yaml"):    _json_to_yaml,
    ("json", "yml"):     _json_to_yaml,
    # csv во всё
    ("csv",  "json"):    _csv_to_json,
    ("csv",  "xml"):     _csv_to_xml,
    # yaml/yml во всё
    ("yaml", "json"):    _yaml_to_json,
    ("yml",  "json"):    _yaml_to_json,
    ("yaml", "xml"):     _yaml_to_xml,
    ("yml",  "xml"):     _yaml_to_xml,
    # docx
    ("docx", "txt"):     _docx_to_txt,
    ("docx", "html"):    _docx_to_html,
    ("docx", "md"):      _docx_to_md,
    ("docx", "pdf"):     _docx_to_pdf,
    ("doc",  "txt"):     _docx_to_txt,
    ("doc",  "html"):    _docx_to_html,
    ("doc",  "pdf"):     _docx_to_pdf,
    # pptx
    ("pptx", "txt"):     _pptx_to_txt,
    ("pptx", "html"):    _pptx_to_html,
    ("pptx", "pdf"):     _pptx_to_pdf,
    ("ppt",  "txt"):     _pptx_to_txt,
    ("ppt",  "pdf"):     _pptx_to_pdf,
    # xlsx/xls
    ("xlsx", "csv"):     _xlsx_to_csv,
    ("xlsx", "json"):    _xlsx_to_json,
    ("xlsx", "html"):    _xlsx_to_html,
    ("xls",  "csv"):     _xlsx_to_csv,
    ("xls",  "json"):    _xlsx_to_json,
    # pdf
    ("pdf",  "txt"):     _pdf_to_txt,
    ("pdf",  "html"):    _pdf_to_html,
    # txt/html/md
    ("txt",  "html"):    _txt_to_html,
    ("txt",  "pdf"):     _txt_to_pdf,
    ("md",   "html"):    _md_to_html,
    ("markdown", "html"):_md_to_html,
    ("html", "txt"):     _html_to_txt,
    ("htm",  "txt"):     _html_to_txt,
}


def list_formats() -> dict:
    # возвращаем что из чего можно сделать
    result: dict[str, list[str]] = {}
    for src_fmt, dst_fmt in sorted(_DISPATCH.keys()):
        result.setdefault(src_fmt, []).append(dst_fmt)
    # картинки умеют всё в себя
    for fmt in IMAGE_FMTS:
        result[fmt] = sorted(IMAGE_FMTS - {fmt})
    result["*видео*"] = ["mp4", "avi", "mkv", "mp3", "wav", "...через ffmpeg"]
    result["*аудио*"] = ["mp3", "wav", "flac", "aac", "ogg", "...через ffmpeg"]
    return result
